from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

bg_dir = os.path.join(os.path.expanduser("~"), "Desktop", "Buildx")
bg = lambda n: os.path.join(bg_dir, f"bg_{n}.png")

GOLD = RGBColor(0xFF, 0xBD, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xDD, 0xDD, 0xDD)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)

def add_bg(slide, n):
    slide.shapes.add_picture(bg(n), Emu(0), Emu(0), prs.slide_width, prs.slide_height)

def tb(slide, l, t, w, h, text, sz=18, col=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = col; p.font.bold = bold; p.font.name = "Calibri"; p.alignment = align

def ml(slide, l, t, w, h, items, sz=15, col=WHITE, sp=4):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = col
        p.font.name = "Calibri"; p.space_after = Pt(sp)

# ============================================================
# SLIDE 1: Title — Week 1 Progress Report
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, 1)

tb(s, 1.0, 2.4, 10, 0.7, "Ann Raksha — Week 1 Progress Report", 30, GOLD, True)

ml(s, 1.0, 3.6, 8, 3.5, [
    "•   TEAM NAME:  Prizzm",
    "",
    "•   TEAM MEMBERS:  Ayush Kushwaha  &  Khushi Pandey",
    "",
    "•   THEME:  Social Impact  |  Sustainability",
    "",
    "•   CURRENT STAGE:  MVP — Functional full-stack application",
    "",
    "•   REPO:  github.com/FrozenLionMax/Ann_Raksha",
], 18, GOLD, 4)

# ============================================================
# SLIDE 2: What We Built This Week
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, 2)

tb(s, 0.8, 1.8, 8, 0.5, "What We Built This Week", 26, GOLD, True)

ml(s, 0.8, 2.5, 6, 4.5, [
    "✅  Complete MERN Stack Application — React 19 frontend + Node.js/Express backend + MongoDB database fully wired together.",
    "",
    "✅  Role-Based Authentication — JWT-secured registration & login for 5 user types: Donor, NGO, Receiver, Volunteer, Admin.",
    "",
    "✅  Real-Time Donation Lifecycle — Full flow from listing surplus food → NGO claiming → volunteer pickup → delivery complete.",
    "",
    "✅  WebSocket Notifications — Socket.IO powered instant alerts when donations are created, claimed, or completed.",
    "",
    "✅  Gamification & Impact Engine — Live dashboard tracking meals served, CO2 prevented, water saved with points & rewards.",
    "",
    "✅  AI Recipe Suggestions — AI-powered module that suggests recipes from surplus food items to reduce waste at source.",
], 14, WHITE, 3)

# Right side - key metrics
tb(s, 8.5, 2.0, 4, 0.5, "Key Deliverables", 20, GOLD, True, PP_ALIGN.CENTER)

tb(s, 8.5, 2.8, 4, 0.7, "20+", 40, GREEN, True, PP_ALIGN.CENTER)
tb(s, 8.5, 3.4, 4, 0.4, "React Components Built", 14, GRAY, False, PP_ALIGN.CENTER)

tb(s, 8.5, 4.0, 4, 0.7, "15+", 40, GREEN, True, PP_ALIGN.CENTER)
tb(s, 8.5, 4.6, 4, 0.4, "API Endpoints Functional", 14, GRAY, False, PP_ALIGN.CENTER)

tb(s, 8.5, 5.2, 4, 0.7, "5", 40, GREEN, True, PP_ALIGN.CENTER)
tb(s, 8.5, 5.8, 4, 0.4, "User Roles Implemented", 14, GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 3: Technical Architecture
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, 3)

tb(s, 0.8, 1.8, 8, 0.5, "Technical Architecture & Stack", 26, GOLD, True)

techs = [
    ("⚛️ Frontend", [
        "React 19 + Vite (fast dev server)",
        "Framer Motion (smooth animations)",
        "Tailwind CSS (responsive design)",
        "Lucide React (icon system)",
        "Dark mode + Glassmorphism UI",
    ]),
    ("🖥️ Backend", [
        "Node.js + Express.js",
        "RESTful API (15+ endpoints)",
        "JWT auth + middleware",
        "Role-based access control",
        "Socket.IO real-time events",
    ]),
    ("🗄️ Database", [
        "MongoDB + Mongoose ODM",
        "User / Donation / Notification",
        "Status lifecycle enums",
        "Deep schema validation",
        "Indexed queries",
    ]),
    ("🧠 Intelligence", [
        "AI recipe suggestions",
        "Impact calculator (CO2/water)",
        "Gamification points engine",
        "Leaderboard system",
        "Real-time dashboard stats",
    ]),
]

for i, (title, items) in enumerate(techs):
    x = 0.5 + i * 3.2
    tb(s, x, 2.5, 3.0, 0.4, title, 18, GOLD, True)
    ml(s, x, 3.0, 3.0, 3.0, [f"•  {item}" for item in items], 13, WHITE, 6)

tb(s, 0.8, 6.0, 12, 0.4, "Architecture:  React UI  →  Axios  →  Express API  →  MongoDB  →  Socket.IO  →  Live Updates", 15, GRAY, False)

# ============================================================
# SLIDE 4: Features Demo
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, 4)

tb(s, 0.8, 1.8, 8, 0.5, "Key Features Implemented", 26, GOLD, True)

features = [
    ("🍽️ Smart Donation Listing", "Donors post surplus food with title, quantity, expiry time, serves count, urgency level & pickup address. Photos supported. One-click publish."),
    ("🔔 Real-Time Notifications", "WebSocket-powered instant alerts — donors get notified when food is claimed, NGOs get pinged when new food is listed nearby. No refresh needed."),
    ("📊 Impact Dashboard", "Role-based live dashboard showing total donations, meals enabled, CO2 prevented (kg), water saved (liters). Updates in real-time after every completed donation."),
    ("🏆 Gamification System", "Points awarded per donation based on weight & urgency. Leaderboards rank top donors. Badges unlock at milestones — makes donating addictive."),
    ("🤖 AI Recipe Engine", "AI suggests creative recipes from surplus food items — helps households reduce waste before it even needs donating."),
    ("🔐 Role-Based Access", "5 distinct user roles (Donor, NGO, Receiver, Volunteer, Admin) each with tailored dashboards, permissions & workflows."),
]

for i, (title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.2
    y = 2.5 + row * 2.4
    tb(s, x, y, 3.9, 0.4, title, 16, GOLD, True)
    ml(s, x, y + 0.45, 3.9, 1.6, [desc], 12, LIGHT, 4)

# ============================================================
# SLIDE 5: Challenges & Learnings
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, 5)

tb(s, 0.8, 1.8, 5.5, 0.5, "Challenges Faced", 22, GOLD, True)

ml(s, 0.8, 2.4, 5.5, 4.5, [
    "⚡  Registration payload mismatch — frontend was sending 'userType' but backend expected 'role'. Caused silent 500 errors. Fixed by aligning API contracts.",
    "",
    "⚡  WebSocket event duplication — Socket.IO was firing duplicate 'status_update' events causing UI flickers. Resolved by implementing proper cleanup on component unmount.",
    "",
    "⚡  Donation tracking bug — After claiming, the Track page was fetching from 'available' donations list, so claimed items vanished from UI. Fixed by building a dedicated GET /donations/:id endpoint.",
    "",
    "⚡  GitHub push protection — Accidentally committed .env with API keys. GitHub blocked the push. Had to rewrite git history to remove secrets and add proper .gitignore.",
], 13, WHITE, 3)

tb(s, 7.2, 1.8, 5.5, 0.5, "Key Learnings", 22, GOLD, True)

ml(s, 7.2, 2.4, 5.5, 4.5, [
    "💡  API contract alignment between frontend and backend should be defined first — prevents hours of debugging payload mismatches.",
    "",
    "💡  Real-time features (WebSockets) need careful lifecycle management — always clean up listeners on unmount to prevent memory leaks.",
    "",
    "💡  Role-based architecture pays off early — designing for multiple user types from day one made the codebase scalable instead of needing painful refactors later.",
    "",
    "💡  Never commit .env files — set up .gitignore before writing the first line of code. GitHub's push protection saved us from a real security incident.",
    "",
    "💡  Gamification is a powerful motivator — even during testing, seeing the impact numbers tick up made the experience feel genuinely rewarding.",
], 13, WHITE, 3)

# ============================================================
# SLIDE 6: Next Week Plan
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, 6)

tb(s, 0.8, 1.8, 8, 0.5, "Week 2 Plan & Next Steps", 26, GOLD, True)

ml(s, 0.8, 2.5, 5.5, 4.5, [
    "🎯  PRIORITY TASKS:",
    "",
    "1.  GPS-Based Live Map — Show nearby donations as pins on a real map. Enable location-based browsing for NGOs.",
    "",
    "2.  Mobile Responsiveness — Ensure the entire UI works flawlessly on phones since donors often post on-the-go.",
    "",
    "3.  Volunteer Dispatch System — Allow volunteers to sign up for pickup routes and get assigned to nearby donations.",
    "",
    "4.  Email Notifications — Backup notification channel via email for users who aren't actively on the platform.",
    "",
    "5.  Admin Analytics Panel — Build comprehensive admin dashboard with city-level stats, user growth, and donation heatmaps.",
], 14, WHITE, 3)

ml(s, 7.2, 2.5, 5.5, 4.5, [
    "🚀  STRETCH GOALS:",
    "",
    "•  Carbon Credit Tracking — Quantify environmental impact into tradeable carbon credits for corporate partners.",
    "",
    "•  Corporate API — OAuth-secured endpoints so enterprises can pull ESG impact data into their own dashboards.",
    "",
    "•  Multi-City Expansion — City-level data isolation and analytics to prepare for scaling beyond one city.",
    "",
    "",
    "📊  SUCCESS METRICS FOR WEEK 2:",
    "",
    "•  Map feature live with real donation pins",
    "•  Mobile-responsive across all pages",
    "•  At least 2 new API endpoints deployed",
], 14, WHITE, 3)

# ============================================================
# SLIDE 7: Thank You
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, 7)

tb(s, 0.8, 1.8, 11, 0.5, "Summary & Links", 26, GOLD, True)

ml(s, 0.8, 2.5, 11, 4, [
    "📦  WHAT WE DELIVERED:  A fully functional MVP with real-time food rescue matching, gamification, AI features, and 5 user roles — all deployed and working.",
    "",
    "🔧  TECH STACK:  React 19 + Vite  |  Node.js + Express  |  MongoDB  |  Socket.IO  |  JWT Auth  |  Framer Motion",
    "",
    "🎯  STAGE:  MVP — Core product complete, ready for user testing and iteration.",
    "",
    "",
    "🔗  LINKS:",
    "",
    "    GitHub:  github.com/FrozenLionMax/Ann_Raksha",
    "",
    "",
    "Thank you, mentors & judges! We'd love your feedback on our approach.",
    "",
    "— Team Prizzm  |  Ayush Kushwaha & Khushi Pandey",
], 16, WHITE, 3)

# SAVE
out = os.path.join(os.path.expanduser("~"), "Desktop", "Prizzm_BuildX26_Week1Update.pptx")
prs.save(out)
print(f"DONE! Saved: {out}")
