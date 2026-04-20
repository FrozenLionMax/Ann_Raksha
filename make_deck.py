from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
# 16:9 widescreen matching the template
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Extract background images from PDF
import fitz
doc = fitz.open(r'C:\Users\Acer\Desktop\Buildx\TeamName_BuildX26 - Google Slides.pdf')
bg_images = []
for i, page in enumerate(doc):
    pix = page.get_pixmap(dpi=150)
    path = os.path.join(os.path.expanduser("~"), "Desktop", "Buildx", f"bg_{i+1}.png")
    pix.save(path)
    bg_images.append(path)

# Colors matching template
GOLD = RGBColor(0xFF, 0xBD, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xDD, 0xDD, 0xDD)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)
EMERALD = RGBColor(0x10, 0xB9, 0x81)

def add_bg_image(slide, img_path):
    slide.shapes.add_picture(img_path, Emu(0), Emu(0), prs.slide_width, prs.slide_height)

def tb(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box

def multi_text(slide, left, top, width, height, items, size=16, color=WHITE, bold=False, spacing=8):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.space_after = Pt(spacing)

# ============================================================
# SLIDE 1: Team & Idea
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_image(s, bg_images[0])

# Idea name in the box area
tb(s, 1.0, 2.4, 10, 0.7, "Ann Raksha — Real-Time Food Rescue Platform", 28, GOLD, True)

# Bullet points below
multi_text(s, 1.0, 3.6, 8, 3.5, [
    "•   ONE-LINE DESCRIPTION:  An AI-powered real-time platform that rescues surplus food from donors and delivers it to NGOs & communities — turning waste into meals.",
    "",
    "•   THEME:  Social Impact  |  Sustainability",
    "",
    "•   TEAM NAME:  Prizzm",
    "",
    "•   TEAM MEMBERS:  Ayush Kushwaha  &  Khushi Pandey",
], 18, GOLD, True, 4)

# ============================================================
# SLIDE 2: Problem Description
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_image(s, bg_images[1])

# Content area — big open space below the title
tb(s, 0.8, 1.8, 7, 0.6, "The Hunger-Waste Paradox", 26, GOLD, True)

multi_text(s, 0.8, 2.5, 7.5, 4.5, [
    "•  India wastes 68.7 million tonnes of food annually — worth ₹92,000 crore — while 190+ million Indians go hungry every day.",
    "",
    "•  Restaurants, hotels, caterers & households throw away perfectly edible food daily because no efficient real-time rescue system exists.",
    "",
    "•  Wasted food in landfills produces methane, contributing to 8-10% of global greenhouse gas emissions.",
    "",
    "•  This is NOT a supply problem — it's a logistics & awareness gap. The food exists; the connection doesn't.",
    "",
    "•  Existing solutions are offline, manual, and fragmented — no platform offers real-time donor-to-NGO matching with impact tracking.",
], 16, WHITE, False, 4)

# Right side stats
tb(s, 9.0, 2.0, 3.5, 0.8, "68.7M Tonnes", 32, GOLD, True, PP_ALIGN.CENTER)
tb(s, 9.0, 2.7, 3.5, 0.5, "food wasted/year in India", 14, GRAY, False, PP_ALIGN.CENTER)

tb(s, 9.0, 3.5, 3.5, 0.8, "190M+ People", 32, RGBColor(0xEF, 0x44, 0x44), True, PP_ALIGN.CENTER)
tb(s, 9.0, 4.2, 3.5, 0.5, "go hungry every day", 14, GRAY, False, PP_ALIGN.CENTER)

tb(s, 9.0, 5.0, 3.5, 0.8, "₹92,000 Cr", 32, EMERALD, True, PP_ALIGN.CENTER)
tb(s, 9.0, 5.7, 3.5, 0.5, "worth of food thrown away", 14, GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 3: Research & Insights
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_image(s, bg_images[2])

# Left column - statistics & research
tb(s, 0.8, 1.8, 6, 0.5, "Key Statistics & Research", 22, GOLD, True)

multi_text(s, 0.8, 2.4, 5.8, 4.5, [
    "•  UNEP Food Waste Index 2024: India ranks among the top 5 food-wasting nations globally.",
    "",
    "•  FAO Report: One-third of all food produced worldwide is lost or wasted — India mirrors this trend at 40%.",
    "",
    "•  FSSAI estimates Indian households waste 50 kg of food per person per year.",
    "",
    "•  Survey of 50+ Delhi NCR restaurants: 78% said they would donate surplus food if pickup was instant and hassle-free.",
    "",
    "•  Interviews with 10+ NGOs: The #1 challenge is unpredictable food supply — they need real-time notifications, not phone calls.",
], 14, WHITE, False, 3)

# Right column - user observations
tb(s, 7.2, 1.8, 5.5, 0.5, "User Observations", 22, GOLD, True)

multi_text(s, 7.2, 2.4, 5.5, 4.5, [
    "•  Donors want a one-click posting system — no paperwork, instant confirmation that their food will reach someone.",
    "",
    "•  NGOs need reliable, real-time alerts when food becomes available nearby — timing is critical for perishables.",
    "",
    "•  Corporates want quantifiable ESG impact metrics (CO2 saved, meals served) for CSR compliance reports.",
    "",
    "•  Volunteers want an easy way to sign up and coordinate pickups without complex logistics.",
    "",
    "•  Everyone wants transparency — proof that their contribution made a real, measurable impact.",
], 14, WHITE, False, 3)

# ============================================================
# SLIDE 4: Proposed Solution
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_image(s, bg_images[3])

tb(s, 0.8, 1.8, 11, 0.6, "Ann Raksha connects surplus food donors with NGOs & communities in real-time.", 20, WHITE, True)

# 3 Key features
tb(s, 0.8, 2.8, 3.6, 0.4, "🔄  Real-Time Matching", 18, GOLD, True)
multi_text(s, 0.8, 3.3, 3.6, 1.8, [
    "Donors list surplus food instantly with quantity, expiry & pickup address. NGOs browse, claim & pick up within the freshness window. WebSocket-powered live notifications keep both parties updated at every step.",
], 14, LIGHT, False, 4)

tb(s, 5.0, 2.8, 3.6, 0.4, "📊  Impact Dashboard", 18, GOLD, True)
multi_text(s, 5.0, 3.3, 3.6, 1.8, [
    "Live tracking of meals served, CO2 emissions prevented, and water saved. Role-based views for donors, NGOs, corporates & admins. Gamification with points, badges & leaderboards drives repeat donations.",
], 14, LIGHT, False, 4)

tb(s, 9.2, 2.8, 3.6, 0.4, "🤖  AI + Lifecycle", 18, GOLD, True)
multi_text(s, 9.2, 3.3, 3.6, 1.8, [
    "AI-powered recipe suggestions from surplus items reduce waste at source. Full donation lifecycle tracking: Available → Matched → Picked Up → Completed. JWT auth with role-based access control.",
], 14, LIGHT, False, 4)

# Flow diagram as text
tb(s, 0.8, 5.5, 12, 0.4, "Donation Flow:", 16, GOLD, True)
tb(s, 0.8, 6.0, 12, 0.5, "📦 Donor Lists Food  →  🔔 NGO Gets Alert  →  🤝 NGO Claims  →  🚚 Volunteer Picks Up  →  ✅ Marked Complete  →  📊 Impact Tracked", 15, WHITE, False)

# ============================================================
# SLIDE 5: Users & Competition
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_image(s, bg_images[4])

# Target users - left
tb(s, 0.8, 1.8, 5.5, 0.5, "Who Will Use It?", 22, GOLD, True)

multi_text(s, 0.8, 2.4, 5.5, 4.5, [
    "🍳  DONORS — Restaurants, hotels, caterers, households & event organizers with surplus edible food. They want a quick, hassle-free way to donate instead of discard.",
    "",
    "🏛️  NGOs & RECEIVERS — Food banks, community kitchens, shelters & underprivileged communities. They need reliable, real-time access to nearby surplus food.",
    "",
    "🏢  CORPORATES — CSR teams needing quantifiable ESG impact data, carbon credit tracking & compliance reports for their sustainability goals.",
    "",
    "🙋  VOLUNTEERS — Individuals who want to help with pickup, delivery & community coordination using a simple, organized platform.",
], 14, WHITE, False, 4)

# Competition - right
tb(s, 7.2, 1.8, 5.5, 0.5, "What Exists Today?", 22, GOLD, True)

multi_text(s, 7.2, 2.4, 5.5, 2, [
    "•  Feeding India (Zomato) — Corporate-only, no real-time matching",
    "•  Robin Hood Army — Volunteer-driven, no tech platform",
    "•  Local NGO WhatsApp groups — Fragmented, no tracking",
    "•  No Indian platform offers live donor-to-NGO matching with impact tracking",
], 14, GRAY, False, 6)

tb(s, 7.2, 4.5, 5.5, 0.5, "How Ann Raksha Differs:", 18, GOLD, True)
multi_text(s, 7.2, 5.0, 5.5, 2.5, [
    "✅  Real-time WebSocket matching (not manual/offline)",
    "✅  Full lifecycle tracking with live status updates",
    "✅  Gamification that drives repeat donor behavior",
    "✅  Quantifiable ESG impact (CO2, water, meals)",
    "✅  AI-powered features (recipe suggestions)",
    "✅  Open to everyone — individuals to corporates",
], 14, EMERALD, False, 4)

# ============================================================
# SLIDE 6: Technologies Used
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_image(s, bg_images[5])

# 4 columns
techs = [
    ("⚛️ Frontend", [
        "React 19 + Vite",
        "Framer Motion animations",
        "Tailwind CSS",
        "Lucide React icons",
        "Dark mode + Glassmorphism",
    ]),
    ("🖥️ Backend", [
        "Node.js + Express.js",
        "RESTful API architecture",
        "JWT authentication",
        "Role-based access control",
        "Socket.IO real-time events",
    ]),
    ("🗄️ Database", [
        "MongoDB + Mongoose ODM",
        "User, Donation, Notification schemas",
        "Deep validation & enums",
        "Indexed queries",
        "Atomic status transitions",
    ]),
    ("🧠 Intelligence", [
        "AI recipe suggestions",
        "Impact calculation engine",
        "CO2 / Water / Meals metrics",
        "Gamification & points",
        "Leaderboard system",
    ]),
]

for i, (title, items) in enumerate(techs):
    x = 0.5 + i * 3.2
    tb(s, x, 1.8, 3.0, 0.4, title, 18, GOLD, True)
    multi_text(s, x, 2.4, 3.0, 3.5, [f"•  {item}" for item in items], 14, WHITE, False, 6)

# Architecture
tb(s, 0.8, 5.6, 12, 0.4, "Architecture Flow:", 16, GOLD, True)
tb(s, 0.8, 6.1, 12, 0.5, "React UI  →  Axios  →  Express API  →  MongoDB  →  Socket.IO  →  Real-Time Updates", 16, WHITE, False)

tb(s, 0.8, 6.8, 12, 0.4, "🔗  github.com/FrozenLionMax/Ann_Raksha", 14, GRAY, False)

# ============================================================
# SLIDE 7: What We Need Feedback On
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_image(s, bg_images[6])

# Assumption
tb(s, 0.8, 1.7, 5.5, 0.5, "🎯  Biggest Assumption", 20, GOLD, True)
multi_text(s, 0.8, 2.3, 5.5, 1.5, [
    "Gamification (points, badges, leaderboards) will be strong enough to drive repeat donor behavior and create a self-sustaining food rescue ecosystem — without needing monetary incentives.",
], 15, WHITE, False, 4)

# Challenge
tb(s, 7.2, 1.7, 5.5, 0.5, "⚡  Biggest Challenge", 20, GOLD, True)
multi_text(s, 7.2, 2.3, 5.5, 1.5, [
    "Building a consistent, reliable supply of food donations so NGOs stay engaged — the classic chicken-and-egg problem. Donors need NGOs to claim, and NGOs need donors to list.",
], 15, WHITE, False, 4)

# Questions
tb(s, 0.8, 4.2, 12, 0.5, "💬  Questions for Mentors & Judges", 22, GOLD, True)

multi_text(s, 0.8, 4.9, 12, 2.5, [
    "1.   Is gamification the right primary driver for sustained donor engagement, or should we explore other incentive models?",
    "",
    "2.   Should we prioritize corporate CSR partnerships (top-down approach) or grassroots community adoption (bottom-up) first?",
    "",
    "3.   Which feature should we build next for maximum impact — mobile app, GPS-based live map, or carbon credit tracking?",
    "",
    "4.   Are we solving the right problem? Is real-time matching the key gap, or is awareness/education more critical?",
], 15, WHITE, False, 3)

# ============================================================
# SAVE
# ============================================================
output = os.path.join(os.path.expanduser("~"), "Desktop", "Prizzm_BuildX26_IdeaSubmission.pptx")
prs.save(output)
print(f"DONE! Saved: {output}")
